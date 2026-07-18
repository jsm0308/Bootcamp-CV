from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUTS = ROOT / "outputs" / "cpp"
PPT_PATH = ROOT / "submission" / "SUBMISSION_WEEK2_RGBD_POINTCLOUD.pptx"

NAVY = RGBColor(20, 27, 35)
INK = RGBColor(28, 34, 41)
MUTED = RGBColor(91, 103, 116)
LINE = RGBColor(216, 222, 228)
CYAN = RGBColor(0, 145, 174)
GREEN = RGBColor(30, 138, 90)
ORANGE = RGBColor(218, 117, 34)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(244, 247, 249)


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Malgun Gothic"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.paragraphs[0].text = text
    for paragraph in frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_title(slide, number, title, subtitle=None):
    add_text(slide, f"0{number}", 0.55, 0.35, 0.55, 0.35, 13, CYAN, True)
    add_text(slide, title, 1.15, 0.25, 11.5, 0.55, 27, INK, True)
    if subtitle:
        add_text(slide, subtitle, 1.15, 0.78, 11.2, 0.35, 11, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.17),
                                  Inches(12.2), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_bullet_list(slide, items, x, y, w, h, size=16, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = 0
    frame.margin_top = 0
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Malgun Gothic"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(9)
        paragraph.text = f"• {paragraph.text}"
    return box


def add_panel(slide, x, y, w, h, fill=PALE, line=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    shape.adjustments[0] = 0.06
    return shape


def add_metric(slide, value, label, x, y, w, color=CYAN):
    add_text(slide, value, x, y, w, 0.45, 22, color, True, PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.45, w, 0.3, 10, MUTED, False, PP_ALIGN.CENTER)


def add_image(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def make_deck():
    PPT_PATH.parent.mkdir(parents=True, exist_ok=True)
    metrics = json.loads((OUTPUTS / "metrics_cpp.json").read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title and verified outcome
    slide = prs.slides.add_slide(blank)
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = NAVY
    add_text(slide, "CV BOOTCAMP · WEEK 2", 0.65, 0.35, 4.2, 0.35, 12,
             RGBColor(99, 211, 226), True)
    add_text(slide, "검증 가능한 RGB-D\n2D → 3D 변환", 0.65, 0.95, 6.0, 1.65, 32, WHITE, True)
    add_text(slide, "OpenCV로 metric point cloud를 만들고\nGoogleTest로 기하·결측·출력을 검증", 0.67, 2.75, 5.5, 0.85,
             17, RGBColor(201, 210, 219))
    add_image(slide, OUTPUTS / "point_cloud_preview_cpp.png", 6.65, 0.55, 6.05, 3.53)
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.72),
                                    Inches(13.333), Inches(1.78))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(238, 243, 246)
    footer.line.fill.background()
    add_metric(slide, "7 / 7", "GoogleTest 통과", 0.75, 6.02, 2.45, GREEN)
    add_metric(slide, f"{metrics['sampled_points']:,}", "생성 point", 3.55, 6.02, 2.45, CYAN)
    add_metric(slide, f"{metrics['valid_depth_ratio'] * 100:.2f}%", "유효 depth", 6.35, 6.02, 2.45, ORANGE)
    add_metric(slide, f"{metrics['max_reprojection_error_px']:.8f} px", "최대 재투영 오차", 9.15, 6.02, 3.1, GREEN)

    # 2. Algorithm
    slide = prs.slides.add_slide(blank)
    add_title(slide, 2, "RGB-D 한 프레임을 metric 3D로 변환",
              "RGB만 보고 깊이를 추정한 것이 아니라, Kinect가 측정한 depth를 사용")
    stages = [
        ("RGB + Depth", "색상과 픽셀별 거리 Z"),
        ("유효값 필터", "0·음수·NaN·Inf 제외"),
        ("Deprojection", "픽셀 좌표를 카메라 좌표로"),
        ("XYZRGB Cloud", "PLY·CSV·시각화 저장"),
    ]
    for i, (name, desc) in enumerate(stages):
        x = 0.65 + i * 3.1
        add_panel(slide, x, 1.55, 2.55, 1.15, WHITE)
        add_text(slide, name, x + 0.15, 1.72, 2.25, 0.35, 17, INK, True, PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.15, 2.12, 2.25, 0.32, 11, MUTED, False, PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, "→", x + 2.6, 1.85, 0.45, 0.45, 22, CYAN, True, PP_ALIGN.CENTER)

    formula = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(3.15),
                                     Inches(5.55), Inches(2.55))
    formula.fill.solid()
    formula.fill.fore_color.rgb = NAVY
    formula.line.fill.background()
    add_text(slide, "Pinhole camera deprojection", 0.95, 3.43, 4.7, 0.35, 13,
             RGBColor(99, 211, 226), True)
    add_text(slide, "X = (u − cx) × Z / fx\nY = (v − cy) × Z / fy\nZ = depth", 0.95, 3.92, 4.7, 1.25,
             22, WHITE, True, PP_ALIGN.LEFT, "Consolas")
    add_text(slide, "(u,v): pixel   ·   (fx,fy,cx,cy): camera intrinsics", 0.95, 5.25, 4.8, 0.25,
             10, RGBColor(193, 202, 211))

    add_text(slide, "왜 depth=0을 버리는가", 6.65, 3.25, 5.55, 0.42, 20, INK, True)
    add_bullet_list(slide, [
        "0은 물체가 카메라 원점에 있다는 뜻이 아니라 측정 실패값",
        "그대로 넣으면 가짜 점들이 (0,0,0) 주변에 쌓여 cloud가 오염",
        "따라서 유효한 Z만 3D 좌표로 변환",
    ], 6.65, 3.83, 5.75, 1.8, 15)

    # 3. Unit tests
    slide = prs.slides.add_slide(blank)
    add_title(slide, 3, "Unit Test: 답을 아는 작은 문제로 검증",
              "Synthetic Golden Case를 먼저 통과시킨 뒤 실제 데이터에 적용")
    add_text(slide, "Golden 3×3 depth map", 0.65, 1.45, 4.2, 0.4, 19, INK, True)
    grid_x, grid_y, cell = 0.85, 2.05, 0.72
    values = [[1, 1, 1], [1, 2, 1], [1, 1, 0]]
    for row in range(3):
        for col in range(3):
            value = values[row][col]
            fill = RGBColor(216, 243, 237) if value == 2 else (RGBColor(244, 221, 215) if value == 0 else PALE)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(grid_x + col * cell), Inches(grid_y + row * cell),
                                           Inches(cell), Inches(cell))
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.color.rgb = WHITE
            add_text(slide, str(value), grid_x + col * cell, grid_y + row * cell,
                     cell, cell, 20, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "K: fx=fy=1, cx=cy=1", 0.85, 4.35, 3.2, 0.35, 13, MUTED)
    add_bullet_list(slide, [
        "중심 (1,1), Z=2 → (0,0,2)",
        "오른쪽 (2,1), Z=1 → (1,0,1)",
        "depth=0 제외 → valid 8, invalid 1",
    ], 0.85, 4.85, 4.4, 1.25, 14)

    tests = [
        ("Golden 3×3 좌표", "PASS"),
        ("결측 depth 필터", "PASS"),
        ("2D→3D→2D round-trip", "PASS"),
        ("잘못된 focal length", "PASS"),
        ("RGB/depth timestamp pairing", "PASS"),
        ("PLY vertex count", "PASS"),
        ("시각화 이미지 생성", "PASS"),
    ]
    add_text(slide, "GoogleTest 결과", 6.0, 1.45, 3.5, 0.4, 19, INK, True)
    for i, (name, status) in enumerate(tests):
        y = 2.0 + i * 0.58
        fill = PALE if i % 2 == 0 else WHITE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(y),
                                       Inches(6.25), Inches(0.48))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
        add_text(slide, name, 6.2, y + 0.04, 4.8, 0.35, 13, INK)
        add_text(slide, status, 11.15, y + 0.04, 0.8, 0.35, 12, GREEN, True, PP_ALIGN.RIGHT)

    # 4. Results
    slide = prs.slides.add_slide(blank)
    add_title(slide, 4, "실제 TUM RGB-D 데이터 결과",
              "Freiburg1 XYZ · 640×480 · Kinect · 첫 번째 RGB/depth association pair")
    add_image(slide, OUTPUTS / "depth_preview_cpp.png", 0.65, 1.45, 3.55, 2.67)
    add_image(slide, OUTPUTS / "projection_preview_cpp.png", 4.35, 1.45, 3.55, 2.67)
    add_image(slide, OUTPUTS / "point_cloud_preview_cpp.png", 8.05, 1.45, 4.65, 2.71)
    add_text(slide, "Depth preview", 0.65, 4.18, 3.55, 0.3, 12, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "3D→2D reprojection", 4.35, 4.18, 3.55, 0.3, 12, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "Point cloud front / top", 8.05, 4.18, 4.65, 0.3, 12, MUTED, True, PP_ALIGN.CENTER)
    y = 5.05
    add_metric(slide, f"{metrics['associated_pairs']}", "RGB-depth pairs", 0.65, y, 2.3, CYAN)
    add_metric(slide, f"{metrics['valid_depth_pixels']:,}", "valid depth pixels", 3.15, y, 2.3, GREEN)
    add_metric(slide, f"{metrics['invalid_depth_pixels']:,}", "invalid depth pixels", 5.65, y, 2.3, ORANGE)
    add_metric(slide, f"{metrics['z_min_m']:.3f}–{metrics['z_max_m']:.3f} m", "depth range", 8.15, y, 2.3, CYAN)
    add_metric(slide, f"{metrics['pair_delta_seconds']:.4f} s", "pair timestamp Δ", 10.65, y, 2.0, GREEN)
    add_text(slide, "재투영 오차는 학습 모델의 정확도가 아니라 구현한 투영·역투영 수식의 일관성을 검증한 값이다.",
             0.65, 6.6, 12.0, 0.35, 11, MUTED, False, PP_ALIGN.CENTER)

    # 5. Conclusion and extension
    slide = prs.slides.add_slide(blank)
    add_title(slide, 5, "결론: RGB-D 기하 파이프라인을 검증 가능한 코드로 구현",
              "단일 프레임 point cloud는 완성, 완전한 3D 복원은 다음 단계")
    columns = [
        (0.65, "완성한 것", GREEN, [
            "OpenCV 16-bit depth 로딩",
            "metric XYZRGB point cloud",
            "결측값·intrinsics 검증",
            "PLY·CSV·3종 시각화",
            "GoogleTest 7/7",
        ]),
        (4.45, "현재 한계", ORANGE, [
            "RGB 단독 depth 예측 아님",
            "가려진 뒷면은 복원 불가",
            "한 프레임의 camera 좌표계",
            "sampling으로 밀도 감소",
            "동적 장면 sync 오차 가능",
        ]),
        (8.25, "Physical AI 확장", CYAN, [
            "pose 기반 multi-view fusion",
            "실물 RGB-D 센서 캡처",
            "object segmentation 결합",
            "grasp 후보와 충돌 영역 분석",
            "Isaac Sim/ROS 2 입력 연결",
        ]),
    ]
    for x, title, color, bullets in columns:
        add_panel(slide, x, 1.55, 3.45, 4.55, WHITE)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.55),
                                     Inches(3.45), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(slide, title, x + 0.25, 1.85, 2.95, 0.45, 19, color, True)
        add_bullet_list(slide, bullets, x + 0.25, 2.55, 2.95, 2.9, 14)

    add_text(slide, "핵심 학습 성과  |  센서가 제공한 depth를 카메라 모델로 해석하고, 작은 정답 문제와 실제 산출물로 코드 정확성을 증명했다.",
             0.65, 6.55, 12.0, 0.4, 12, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "Data: TUM RGB-D Dataset · Implementation: C++17 / OpenCV 4.12 / GoogleTest 1.15.2",
             0.65, 7.0, 12.0, 0.25, 9, MUTED, False, PP_ALIGN.CENTER)

    prs.core_properties.title = "Unit Test 구성 및 RGB-D 2D→3D 변환 알고리즘 설계"
    prs.core_properties.subject = "CV Bootcamp Week 2"
    prs.core_properties.author = "장승민"
    prs.save(PPT_PATH)
    print(PPT_PATH)


if __name__ == "__main__":
    make_deck()
