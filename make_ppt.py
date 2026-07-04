"""
PPT Generator — reusable across weeks.

Usage:
  python make_ppt.py              # build week 1 (default)
  python make_ppt.py --week 2     # build week 2
  python make_ppt.py -o my.pptx   # custom output path
"""
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ──────────────────────────────────────────
NAVY_DARK  = "2c3e50"
NAVY       = "34495e"
WHITE      = "FFFFFF"
BLACK      = "333333"
GRAY       = "555555"
LIGHT_GRAY = "999999"
RED        = "e74c3c"
GREEN      = "27ae60"
BLUE       = "2980b9"

# ── Layout constants ────────────────────────────────────────
W, H = 13.333, 7.5  # 16:9

# ═══════════════════════════════════════════════════════════════
#  Builder primitives (reusable across all weeks)
# ═══════════════════════════════════════════════════════════════

def _init():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    return prs, lambda: prs.slides.add_slide(prs.slide_layouts[6])

def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def _para(tf, text, size=20, bold=False, color=BLACK, align=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
    p.font.color.rgb = RGBColor(r, g, b)
    if align:
        p.alignment = align
    return p

def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
    fill.fore_color.rgb = RGBColor(r, g, b)

def _img(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        print(f"  [WARN] image not found: {path}")
        return None
    if width:
        return slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
    if height:
        return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                        height=Inches(height))
    return slide.shapes.add_picture(path, Inches(left), Inches(top))

# ═══════════════════════════════════════════════════════════════
#  Slide builders (content passed in, not hardcoded)
# ═══════════════════════════════════════════════════════════════

def cover_slide(new_slide, title, subtitle, meta, accent=""):
    """
    Dark navy cover page.
    title:   main heading (str)
    subtitle: sub-heading (str)
    meta:    bottom text lines (list of str)
    accent:  optional accent text below subtitle (str)
    """
    slide = new_slide()
    _bg(slide, NAVY_DARK)
    tf = _box(slide, 1, 1.5, 11, 4)
    _para(tf, title, size=44, bold=True, color=WHITE)
    _para(tf, subtitle, size=26, color="BBBBBB")
    if accent:
        _para(tf, "", size=10)
        _para(tf, accent, size=16, color=RED)
    _para(tf, "", size=12)
    for line in meta:
        _para(tf, line, size=16, color=LIGHT_GRAY)

def section_slide(new_slide, title, body):
    """
    Navy section divider.
    title: large heading (str)
    body:  subtitle lines (list of str)
    """
    slide = new_slide()
    _bg(slide, NAVY)
    tf = _box(slide, 1, 2.0, 11, 3.5)
    _para(tf, title, size=38, bold=True, color=WHITE)
    for line in body:
        _para(tf, line, size=20, color="BBBBBB")

def content_slide(new_slide, title, content, images=None,
                  code_block=None, note=None):
    """
    White content slide with navy header bar.
    title:      page heading (str)
    content:    list of (text, size, bold, color) tuples
    images:     optional list of (path, left, top, width) tuples
    code_block: optional list of (code_line, color) tuples
    note:       optional footnote (str)
    """
    slide = new_slide()
    # Navy top bar
    from pptx.util import Inches, Pt
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(W), Inches(1.1))  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5e)
    bar.line.fill.background()
    # Title on bar
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.15), Inches(11), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Body text
    tf = _box(slide, 1.5, 1.5, 10, 4.5)
    for item in content:
        if len(item) == 4:
            text, size, bold, color = item
            _para(tf, text, size=size, bold=bold, color=color)
        elif len(item) == 2:
            text, size = item
            _para(tf, text, size=size)

    # Code block (right column)
    if code_block:
        cf = _box(slide, 7.0, 1.5, 5, 5.0)
        _para(cf, "Core Code", size=16, bold=True, color=NAVY)
        for line, clr in code_block:
            _para(cf, line, size=14, color=clr if clr else GRAY)

    # Images
    if images:
        for path, left, top, width in images:
            _img(slide, path, left, top, width=width)

    # Footnote
    if note:
        _para(_box(slide, 1.5, 6.5, 10, 0.5),
              note, size=14, color=GRAY)

def two_column_slide(new_slide, title, left_items, right_items, note=None):
    """
    Two-column content slide. Good for before/after or comparison.
    left_items:  list of (text, size, bold, color)
    right_items: list of (text, size, bold, color)
    """
    slide = new_slide()
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5e)
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.15), Inches(11), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Left
    tf = _box(slide, 1.0, 1.5, 5.2, 5.0)
    for item in left_items:
        text, size, bold, color = item
        _para(tf, text, size=size, bold=bold, color=color)

    # Right
    tf = _box(slide, 7.0, 1.5, 5.2, 5.0)
    for item in right_items:
        text, size, bold, color = item
        _para(tf, text, size=size, bold=bold, color=color)

    if note:
        _para(_box(slide, 1.5, 6.5, 10, 0.5),
              note, size=14, color=GRAY)

def image_grid_slide(new_slide, title, rows, note=None):
    """
    Grid of image pairs (e.g. original vs processed).
    rows: list of dicts with keys: label_left, path_left, label_right, path_right
    """
    slide = new_slide()
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5e)
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.15), Inches(11), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    col_w = 2.2
    for i, row in enumerate(rows):
        y = 1.8 + i * 2.5
        x1, x2 = 2.0, 7.5
        # Labels
        _para(_box(slide, x1, y - 0.3, col_w, 0.4),
              row["label_left"], size=14, color=GRAY, align=PP_ALIGN.CENTER)
        _para(_box(slide, x2, y - 0.3, col_w, 0.4),
              row["label_right"], size=14, color=GRAY, align=PP_ALIGN.CENTER)
        _img(slide, row["path_left"], x1, y, width=col_w)
        _img(slide, row["path_right"], x2, y, width=col_w)

    if note:
        _para(_box(slide, 1.5, 6.5, 10, 0.5), note, size=14, color=GRAY)

def checklist_slide(new_slide, title, items, lessons, bg=NAVY_DARK):
    """
    Final slide: checklist + lessons on dark navy.
    items:   list of checklist strings
    lessons: list of (line, is_title) tuples
    """
    slide = new_slide()
    _bg(slide, bg)
    # Checklist (left)
    _para(_box(slide, 1.2, 0.6, 5, 0.8),
          "Checklist", size=30, bold=True, color=WHITE)
    tf = _box(slide, 1.2, 1.5, 5, 5.0)
    # Table header
    _para(tf, "    Status  Item", size=16, bold=True, color="BBBBBB")
    for item in items:
        _para(tf, f"    [O]    {item}", size=16, color=GREEN)

    # Lessons (right)
    _para(_box(slide, 7.0, 0.6, 5, 0.8),
          "Lessons Learned", size=30, bold=True, color=WHITE)
    tf = _box(slide, 7.0, 1.5, 5.5, 5.0)
    for text, is_title in lessons:
        if text == "":
            _para(tf, "", size=6)
        else:
            _para(tf, text, size=16 if is_title else 14,
                  bold=is_title, color=WHITE if is_title else "AAAAAA")

# ═══════════════════════════════════════════════════════════════
#  Week-specific build functions
# ═══════════════════════════════════════════════════════════════

def build_week1(new_slide):
    """Assemble all slides for Week 1 assignment."""

    # ── Slide 1: Cover ─────────────────────────────────────
    cover_slide(new_slide,
        title="1차 과제 결과 보고서",
        subtitle="Git을 활용한 코드 관리 및 픽셀 단위 이미지 처리 실습",
        meta=["코멘토 직무부트캠프 Computer Vision",
              "베어곰 멘토  |  2026.07"])

    # ── Slide 2: Overview (Background + Topic + Requests) ──
    content_slide(new_slide,
        title="과제 개요",
        content=[
            ("[배경]", 22, True, NAVY),
            ("AI 기반 신제품 조직 신설, Computer Vision 기술로 사용자 편의성 증대 목표.", 18, False, BLACK),
            ("효율적 코드 관리와 픽셀 단위 이미지 처리는 CV 프로젝트 기초 역량.", 18, False, BLACK),
            ("", 8, False, BLACK),
            ("[주제]   Git을 활용한 코드 관리 및 픽셀 단위 이미지 처리 실습", 22, True, NAVY),
            ("", 8, False, BLACK),
            ("[요청내용]", 22, True, NAVY),
            ("1. Git 저장소 구성 및 실습", 18, False, BLACK),
            ("    - GitHub 저장소 생성, 로컬 연동, 브랜치 관리", 16, False, GRAY),
            ("2. 픽셀 단위 이미지 처리 코드 작성", 18, False, BLACK),
            ("    - OpenCV로 이미지 로드, HSV 변환, 특정 색상 필터링", 16, False, GRAY),
            ("3. Git 코드 관리 및 제출", 18, False, BLACK),
            ("    - commit, push, PR 생성, merge", 16, False, GRAY),
        ])

    # ── Slide 3: Git setup + branch ────────────────────────
    content_slide(new_slide,
        title="Git 저장소 구성 & Branch 실습",
        content=[
            ("1.1 GitHub 저장소 생성", 22, True, NAVY),
            ('    "초기화 안 함" 상태로 생성 (No README / No .gitignore)', 16, False, GRAY),
            ("    저장소 URL 복사 → 로컬 환경에 연동", 16, False, GRAY),
            ("", 8, False, BLACK),
            ("1.2 로컬 환경 Git 연동", 22, True, NAVY),
            ("    > git init", 16, False, BLUE),
            ("    > git clone [저장소 URL]", 16, False, BLUE),
            ("    > git pull origin main", 16, False, BLUE),
            ("", 8, False, BLACK),
            ("2. Branch 및 Commit 실습", 22, True, NAVY),
            ("    > git branch feature/image-processing", 16, False, BLUE),
            ("    > git checkout feature/image-processing", 16, False, BLUE),
            ("", 8, False, BLACK),
            ("[결과] jsm0308/Bootcamp-CV  +  feature/image-processing 브랜치", 18, True, GREEN),
        ])

    # ── Slide 4: Image processing code ─────────────────────
    repo_img = "week 1/jsm0308-Bootcamp-CV.png"
    images = []
    if os.path.exists(repo_img):
        images.append((repo_img, 7.8, 2.2, 4.5))

    content_slide(new_slide,
        title="이미지 처리 코드 작성",
        content=[
            ("OpenCV로 이미지 로드 → HSV 색상 공간 변환 → 빨간색 필터링", 20, True, NAVY),
            ("", 6, False, BLACK),
            ("사용 함수:", 18, False, BLACK),
            ("  cv2.imread()         이미지 로드", 16, False, GRAY),
            ("  cv2.cvtColor()       BGR → HSV 변환", 16, False, GRAY),
            ("  cv2.inRange()        빨간색 범위 마스킹 (2개)", 16, False, GRAY),
            ("  cv2.bitwise_or()     두 마스크 합치기", 16, False, RED),
            ("  cv2.bitwise_and()    원본 × 마스크 → 결과", 16, False, GRAY),
            ("  cv2.imwrite()        결과 저장", 16, False, GRAY),
            ("", 6, False, BLACK),
            ("[주의] mask1 + mask2 → cv2.bitwise_or 로 교체", 18, False, RED),
            ("  uint8에서 255+255=510 overflow 방지", 14, False, GRAY),
        ],
        code_block=[
            ("import cv2, numpy as np", GRAY),
            ("image = cv2.imread('sample.jpg')", GRAY),
            ("hsv = cv2.cvtColor(image,", GRAY),
            ("          cv2.COLOR_BGR2HSV)", GRAY),
            ("", None),
            ("lower_red1 = [0, 120, 70]", RED),
            ("upper_red1 = [10, 255, 255]", RED),
            ("lower_red2 = [170, 120, 70]", RED),
            ("upper_red2 = [180, 255, 255]", RED),
            ("", None),
            ("mask1 = cv2.inRange(hsv, ...)", GRAY),
            ("mask2 = cv2.inRange(hsv, ...)", GRAY),
            ("mask = cv2.bitwise_or(m1, m2)", GRAY),
            ("result = cv2.bitwise_and(", GRAY),
            ("    image,image, mask=mask)", GRAY),
            ("cv2.imwrite(", GRAY),
            ("    'red_filtered.jpg',", GRAY),
            ("    result)", GRAY),
        ],
        images=images,
        note="POC: red.py는 추가 라이브러리 없이 OpenCV + NumPy만으로 구현")

    # ── Slide 5: Git submit + PR ────────────────────────────
    content_slide(new_slide,
        title="Git 최종 제출 & Pull Request",
        content=[
            ("4. Git 최종 코드 제출", 22, True, NAVY),
            ("    > git status                  # 변경 사항 확인", 16, False, GRAY),
            ("    > git add .                   # 변경된 파일 추가", 16, False, GRAY),
            ('    > git commit -m "..message.."  # 커밋', 16, False, GRAY),
            ("    > git push origin [...]       # 원격 저장소 업로드", 16, False, GRAY),
            ("", 8, False, BLACK),
            ("5. Pull Request & Merge", 22, True, NAVY),
            ("    GitHub에서 PR 생성 → 리뷰어 추가 → 코드 리뷰", 16, False, GRAY),
            ("    > git checkout main           # main 브랜치로 이동", 16, False, GRAY),
            ("    > git merge [브랜치 이름]     # 변경 내용 병합", 16, False, GRAY),
            ("    > git push origin main        # 최종 푸시", 16, False, GRAY),
            ("", 8, False, BLACK),
            ("[결과] PR #1 생성 → main Fast-forward merge 완료", 18, True, GREEN),
            ("        https://github.com/jsm0308/Bootcamp-CV/pull/1", 14, False, BLUE),
        ])

    # ── Slide 6: HuggingFace preprocessing ──────────────────
    content_slide(new_slide,
        title="추가 요청: HuggingFace 전처리",
        content=[
            ("HuggingFace ethz/food101", 20, True, NAVY),
            ("  75,750장 / 101종 / 평균 502x476 px", 16, False, GRAY),
            ("", 6, False, BLACK),
            ("[기본 문제]", 22, True, RED),
            ("  1) 크기 조정      → cv2.resize() 224 X 224", 16, False, GRAY),
            ("  2) Grayscale      → cv2.cvtColor(..., GRAY)", 16, False, GRAY),
            ("  3) Normalize      → / 255.0  [0, 1]", 16, False, GRAY),
            ("  4) Blur 필터     → cv2.GaussianBlur(5x5)", 16, False, GRAY),
            ("  5) 데이터 증강   → 좌우반전, 15도회전, 밝기+30", 16, False, GRAY),
            ("", 6, False, BLACK),
            ("[심화 문제]  이상치 탐지", 22, True, RED),
            ("  - 너무 어두운 이미지: np.mean(gray) < 30 → skip", 16, False, GRAY),
            ("  - 객체 작은 이미지: OTSU threshold → contour 면적 < 10% → skip", 16, False, GRAY),
            ("", 6, False, BLACK),
            ("의존성: OpenCV + NumPy + datasets + pandas 만 사용", 18, True, GREEN),
        ],
        note="image_preprocessing.py 전체 파이프라인 → 5개 샘플 × 4종 증강 = preprocessed_samples/ 20장")

    # ── Slide 7: Before/After + Repo screenshot ─────────────
    image_grid_slide(new_slide,
        title="Before / After",
        rows=[
            {
                "label_left": "Original",
                "path_left": "week 1/git/sample.jpg",
                "label_right": "Red Filtered",
                "path_right": "week 1/git/red_filtered.jpg",
            },
            {
                "label_left": "Raw Food101",
                "path_left": "week 1/huggingface/data/01_beignets_raw.jpg",
                "label_right": "Preprocessed",
                "path_right": "week 1/huggingface/preprocessed_samples/01_beignets_preprocessed.jpg",
            },
        ],
        note="left: original images | right: after processing")

    # GitHub repo screenshot
    repo_img = "week 1/jsm0308-Bootcamp-CV.png"
    if os.path.exists(repo_img):
        content_slide(new_slide,
            title="GitHub Repository",
            content=[
                ("저장소: jsm0308/Bootcamp-CV", 20, True, NAVY),
                ("git/  → 핵심 과제: red.py (빨간색 검출)", 16, False, GRAY),
                ("huggingface/ → 추가 과제: dataset.py, image_preprocessing.py", 16, False, GRAY),
            ],
            images=[(repo_img, 7.5, 1.8, 5.0)])

    # ── Slide 8: Checklist + Lessons ────────────────────────
    checklist_slide(new_slide,
        title="",
        items=[
            "Git 저장소 구성",
            "Branch / PR / Merge",
            "red.py 빨간색 검출",
            "dataset.py Food101 + pandas",
            "image_preprocessing.py",
            "Resize 224x224",
            "Grayscale + Normalize",
            "GaussianBlur 노이즈 제거",
            "증강 3종 (flip, rotate, bright)",
            "이상치 탐지 2종",
            "전처리 결과물 20장",
            "PPT 보고서",
        ],
        lessons=[
            ("1. HSV가 RGB보다 실용적", True),
            ("   조명 변화에 강하고, 색상 개념에 가깝다", False),
            ("", False),
            ("2. cv2.bitwise_or vs + 연산자", True),
            ("   255+255=510 overflow. 비트연산이 정답.", False),
            ("", False),
            ("3. pandas EDA 먼저, 전처리는 그 다음", True),
            ("   데이터 분포 먼저 파악하고 전략 수립", False),
            ("", False),
            ("4. Git branch는 혼자여도 의미있다", True),
            ("   실수 복구 쉽고, 작업 단위가 명확", False),
            ("", False),
            ("5. OpenCV+NumPy만으로 충분", True),
            ("   추가 의존성 없이 기본 함수 조합으로 구현", False),
        ])

# ═══════════════════════════════════════════════════════════════
#  Main entry point
# ═══════════════════════════════════════════════════════════════

BUILDERS = {
    1: build_week1,
}

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate assignment PPT")
    parser.add_argument("--week", type=int, default=1, help="Which week (1-4)")
    parser.add_argument("-o", "--output", default="presentation.pptx")
    args = parser.parse_args()

    builder = BUILDERS.get(args.week)
    if builder is None:
        print(f"ERROR: week {args.week} not implemented yet.")
        exit(1)

    prs, new_slide = _init()
    builder(new_slide)
    prs.save(args.output)
    print(f"DONE: {args.output} (week {args.week})")
